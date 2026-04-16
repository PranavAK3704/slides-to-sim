"""
ppt_extractor.py
================
Extracts click-sequence steps from a PPTX file for PCT detection.
Ported from valmo-ops/scripts — now runs centrally in the admin backend.

Flow:
  load PPTX
  → segment into processes (by title-slide boundaries)
  → for each process: find red-outlined shapes → extract nearest text labels
  → fallback: parse bottom instruction text "Go to X(1). Click Y(2)."
  → return [{process_name, steps: [{order, elementText, urlPattern}]}]

urlPattern is inferred from the first word of the tab name (lowercase).
Admin can correct it in the UI.
"""

import re
import logging
import tempfile
import requests
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# ── Red-box detection ─────────────────────────────────────────────────────────

RED_MIN_R, RED_MAX_G, RED_MAX_B = 150, 100, 100


def _is_red_shape(shape) -> bool:
    try:
        if not hasattr(shape, 'line'):
            return False
        rgb = shape.line.color.rgb
        if not rgb:
            return False
        r, g, b = rgb
        return r >= RED_MIN_R and g < RED_MAX_G and b < RED_MAX_B
    except Exception:
        return False


def _shape_text(shape) -> Optional[str]:
    try:
        t = shape.text.strip() if hasattr(shape, 'text') else None
        return t if t else None
    except Exception:
        return None


def _nearest_text(slide, ref_shape, max_dist: int = 2_000_000) -> Optional[str]:
    best, best_dist = None, float('inf')
    for shape in slide.shapes:
        if shape is ref_shape:
            continue
        t = _shape_text(shape)
        if not t or len(t) > 100:
            continue
        dist = abs(shape.left - ref_shape.left) + abs(shape.top - ref_shape.top)
        if dist < best_dist and dist < max_dist:
            best, best_dist = t, dist
    return best


# ── Instruction text parsing ──────────────────────────────────────────────────

_INSTRUCTION_RE = re.compile(
    r'(?:go to|click on|click|select|tap|navigate to)\s+([A-Za-z][A-Za-z0-9 /\-]+?)\s*\(\d+\)',
    re.IGNORECASE,
)
_ACTION_VERBS = {'navigate to', 'go to', 'click on', 'click', 'then click',
                 'select', 'open', 'choose', 'tap'}


def _clean_label(raw: str) -> str:
    label = raw.replace('\n', ' ').replace('\r', ' ')
    label = ' '.join(label.split())
    # Strip leading action verbs
    for verb in sorted(_ACTION_VERBS, key=len, reverse=True):
        if label.lower().startswith(verb):
            label = label[len(verb):].strip()
            break
    # Strip trailing "(1)" etc.
    label = re.sub(r'\s*\(\d+\)\s*$', '', label).strip()
    return label


def _instruction_text(slide) -> Optional[str]:
    """Find the long instruction block at the bottom of a slide."""
    for shape in slide.shapes:
        t = _shape_text(shape)
        if t and len(t) > 50 and any(v in t.lower() for v in ('go to', 'click', 'select', 'tap')):
            return t
    return None


def _parse_instruction(text: str) -> list[str]:
    """Extract ordered tab labels from 'Go to X(1). Click Y(2).' pattern."""
    matches = _INSTRUCTION_RE.findall(text or '')
    seen, result = set(), []
    for m in matches:
        label = _clean_label(m)
        if label and label not in seen:
            seen.add(label)
            result.append(label)
    return result


# ── URL pattern inference ─────────────────────────────────────────────────────

def _infer_url_pattern(label: str) -> str:
    """
    Best-effort: take the most specific word from the label.
    e.g. 'RTO Manifest'  → 'rto'
         'Misroute Bagging' → 'misroute'
         'Forward Dispatch' → 'forward'
    """
    first = label.strip().split()[0].lower() if label.strip() else ''
    # Strip common generic words
    generics = {'all', 'the', 'new', 'create', 'open', 'go', 'click', 'select', 'view'}
    words = [w.lower() for w in label.split() if w.lower() not in generics]
    return words[0] if words else first


# ── Process segmentation (by title-slide boundaries) ─────────────────────────

_TITLE_KEYWORDS = ('title', 'divider', 'section')


def _is_title_slide(slide) -> bool:
    layout_name = getattr(slide.slide_layout, 'name', '').lower()
    if any(k in layout_name for k in _TITLE_KEYWORDS):
        return True
    # Fallback: single large text, no screenshots
    texts = [s.text.strip() for s in slide.shapes if hasattr(s, 'text') and s.text.strip()]
    images = [s for s in slide.shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
    return len(texts) <= 2 and len(images) == 0


def _segment(prs) -> list[dict]:
    """Split deck into [{name, slide_indices}] segments."""
    segments, current = [], None
    for i, slide in enumerate(prs.slides):
        if _is_title_slide(slide):
            if current:
                segments.append(current)
            # Extract process name from title slide
            name = ''
            for shape in slide.shapes:
                t = _shape_text(shape)
                if t and len(t) > 3:
                    name = t
                    break
            current = {'name': _clean_label(name) or f'Process {len(segments)+1}', 'slides': []}
        elif current is not None:
            current['slides'].append(i)

    if current and current['slides']:
        segments.append(current)

    # If no title slides found, treat entire deck as one process
    if not segments:
        segments = [{'name': 'Process 1', 'slides': list(range(len(prs.slides)))}]

    return segments


# ── Main extraction ───────────────────────────────────────────────────────────

def extract_from_pptx(ppt_path: str) -> list[dict]:
    """
    Returns list of:
      {
        process_name: str,
        steps: [{order: int, elementText: str, urlPattern: str}]
      }
    """
    from pptx import Presentation  # lazy import so module loads even without pptx
    prs = Presentation(ppt_path)
    segments = _segment(prs)
    results = []

    for seg in segments:
        labels = []
        for idx in seg['slides']:
            if idx >= len(prs.slides):
                continue
            slide = prs.slides[idx]

            # Try instruction text first (most reliable)
            instr = _instruction_text(slide)
            if instr:
                parsed = _parse_instruction(instr)
                if parsed:
                    for lbl in parsed:
                        if lbl not in labels:
                            labels.append(lbl)
                    continue  # skip red-box scan for this slide

            # Fallback: red-outlined shapes
            for shape in slide.shapes:
                if _is_red_shape(shape):
                    nearby = _nearest_text(slide, shape)
                    if nearby:
                        lbl = _clean_label(nearby)
                        if lbl and lbl not in labels:
                            labels.append(lbl)

        if not labels:
            logger.warning(f"No steps found for '{seg['name']}'")
            continue

        steps = [
            {'order': i + 1, 'elementText': lbl, 'urlPattern': _infer_url_pattern(lbl)}
            for i, lbl in enumerate(labels)
        ]
        results.append({'process_name': seg['name'], 'steps': steps})

    return results


# ── Drive download helper ─────────────────────────────────────────────────────

def download_from_url(url: str) -> str:
    """
    Download a PPTX from a Google Slides or Google Drive URL.
    Returns the temp file path.

    Handles:
    - Google Slides: https://docs.google.com/presentation/d/{ID}/...
      → exports as PPTX via /export/pptx
    - Google Drive:  https://drive.google.com/file/d/{ID}/...
      → direct download
    """
    # ── Google Slides export ──
    slides_match = re.search(r'docs\.google\.com/presentation/d/([a-zA-Z0-9_-]+)', url)
    if slides_match:
        pres_id  = slides_match.group(1)
        download_url = f"https://docs.google.com/presentation/d/{pres_id}/export/pptx"
        session  = requests.Session()
        resp     = session.get(download_url, stream=True, allow_redirects=True)
        if not resp.ok:
            raise ValueError(f"Google Slides export failed ({resp.status_code}) — make sure the deck is shared as 'Anyone with link'")
        return _write_tmp(resp)

    # ── Google Drive download ──
    match = re.search(r'/d/([a-zA-Z0-9_-]+)', url) or \
            re.search(r'[?&]id=([a-zA-Z0-9_-]+)', url)
    if not match:
        raise ValueError(f"Could not extract file ID from URL: {url}")

    file_id      = match.group(1)
    download_url = f"https://drive.google.com/uc?export=download&id={file_id}"
    session      = requests.Session()
    resp         = session.get(download_url, stream=True)
    # Handle virus-scan confirmation cookie for large files
    for k, v in resp.cookies.items():
        if k.startswith('download_warning'):
            resp = session.get(f"{download_url}&confirm={v}", stream=True)
            break
    return _write_tmp(resp)


def _write_tmp(resp) -> str:
    """Write a streaming response to a temp PPTX file, return path."""
    tmp = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
    for chunk in resp.iter_content(32768):
        if chunk:
            tmp.write(chunk)
    tmp.flush()
    tmp.close()

    size = Path(tmp.name).stat().st_size
    if size < 1000:
        with open(tmp.name, 'rb') as f:
            if b'<!DOCTYPE' in f.read(200):
                raise ValueError(
                    "Got an HTML page instead of a PPTX — "
                    "file must be shared as 'Anyone with link' (View)"
                )
    return tmp.name


# Keep old name as alias so existing callers don't break
download_from_drive = download_from_url
